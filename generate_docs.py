import os
from docx import Document
from docx.shared import Inches, Pt
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt

# Project Details
PROJECT_NAME = "Servu"
MOTTO = "Build. Link. Grow. Empowering Bulawayo’s Future."
MEMBERS = [
    {
        "name": "Mphathisi Ndlovu",
        "reg": "N02529052A",
        "phone": "0787 146 103",
        "faculty": "Computer Science"
    },
    {
        "name": "Anesu Hove",
        "reg": "N02530834W",
        "phone": "0780 884 195",
        "faculty": "Quantity Surveying"
    },
    {
        "name": "Siyabonga Khumalo",
        "reg": "N02534243G",
        "phone": "0788 532 354",
        "faculty": "Computer Science"
    },
    {
        "name": "Chikomborero I Zvinoira",
        "reg": "N02532066J",
        "phone": "0781 013 061",
        "faculty": "Chemical Engineering"
    }
]

WHAT_IS_SERVU = (
    "Servu (formerly FixBulawayo) is a trust-first digital marketplace designed to connect Bulawayo's "
    "youth with local residents and SMEs for micro-tasks and professional services. Built on the "
    "principles of reliability and digital visibility, it provides a seamless bridge between talent "
    "and opportunity."
)

WHO_IS_IT_FOR = (
    "1. Verified Youth: Skill-holders looking for employment, digital reputation, and a professional workspace.\n"
    "2. Local SMEs: Small businesses needing reliable, verified talent for tasks without the overhead of traditional hiring.\n"
    "3. Residents: Homeowners seeking trusted service providers for domestic or technical micro-tasks."
)

WHY_IMPORTANT = (
    "• Reducing Unemployment: Directly addresses the youth unemployment crisis by enabling the gig economy.\n"
    "• Trust Ecosystem: Features an innovative escrow-backed payment system and a rating-based trust model.\n"
    "• Digital Visibility: Every task completed builds a 'Reputation CV', giving youth a verifiable work history.\n"
    "• Community-Centric: Tailored specifically for the Bulawayo context, with a 'data-light' philosophy for accessibility."
)

LOGO_PATH = "/home/mphatic/Desktop/fixforme/logo_nust_png.png"

def create_word():
    doc = Document()
    
    # Title Page
    if os.path.exists(LOGO_PATH):
        doc.add_picture(LOGO_PATH, width=Inches(1.5))
        last_paragraph = doc.paragraphs[-1]
        last_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER

    title = doc.add_heading(PROJECT_NAME, 0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    subtitle = doc.add_paragraph(MOTTO)
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    doc.add_page_break()
    
    # Group Members
    doc.add_heading('Group R: Team Details', level=1)
    table = doc.add_table(rows=1, cols=4)
    table.style = 'Table Grid'
    hdr_cells = table.rows[0].cells
    hdr_cells[0].text = 'Name'
    hdr_cells[1].text = 'Reg Number'
    hdr_cells[2].text = 'Faculty'
    hdr_cells[3].text = 'Phone'
    
    for member in MEMBERS:
        row_cells = table.add_row().cells
        row_cells[0].text = member['name']
        row_cells[1].text = member['reg']
        row_cells[2].text = member['faculty']
        row_cells[3].text = member['phone']
    
    doc.add_paragraph() # Spacer
    
    # Project Info
    doc.add_heading('What is Servu?', level=1)
    doc.add_paragraph(WHAT_IS_SERVU)
    
    doc.add_heading('Who is it built for?', level=1)
    doc.add_paragraph(WHO_IS_IT_FOR)
    
    doc.add_heading('Why is it Important?', level=1)
    doc.add_paragraph(WHY_IMPORTANT)
    
    doc.save('Servu_Submission.docx')
    print("Word document created: Servu_Submission.docx")

def create_pptx():
    prs = Presentation()
    
    # Slide 1: Title
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = PROJECT_NAME
    subtitle.text = MOTTO
    
    # Slide 2: Team
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Group R: The Team"
    body = slide.placeholders[1]
    tf = body.text_frame
    for m in MEMBERS:
        p = tf.add_paragraph()
        p.text = f"{m['name']} ({m['reg']}) - {m['faculty']}"
        p.level = 0
        
    # Slide 3: What is Servu?
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "The Innovation: Servu"
    body = slide.placeholders[1]
    body.text = WHAT_IS_SERVU
    
    # Slide 4: Target Audience
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Who is it built for?"
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "Empowering key stakeholders:"
    highlights = ["Verified Youth Professionals", "Local SMEs & Businesses", "Community Residents"]
    for h in highlights:
        p = tf.add_paragraph()
        p.text = h
        p.level = 1

    # Slide 5: Importance & Impact
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Why it Matters"
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "Driving change in Bulawayo:"
    impacts = [
        "Digital Visibility & Reputation CVs",
        "Trust through Escrow Payments",
        "Sustainable Youth Employment",
        "Safe Peer-to-Peer Service Economy"
    ]
    for i in impacts:
        p = tf.add_paragraph()
        p.text = i
        p.level = 1
        
    prs.save('Servu_Pitch.pptx')
    print("PowerPoint presentation created: Servu_Pitch.pptx")

if __name__ == "__main__":
    create_word()
    create_pptx()
